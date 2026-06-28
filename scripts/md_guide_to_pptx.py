#!/usr/bin/env python3
"""수련회 사용 가이드 마크다운 → PowerPoint 변환 (일회성 유틸)."""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
GUIDE_MD = ROOT / "docs" / "retreat-leader-guide.md"
IMG_DIR = ROOT / "app" / "retreat" / "static" / "retreat" / "guides"
DEFAULT_OUT = ROOT / "church-retreat-platform-guide.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_W = SLIDE_W - MARGIN * 2


def strip_md(text: str) -> str:
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    return text.strip()


def parse_blocks(md: str) -> list[dict]:
    blocks: list[dict] = []
    lines = md.splitlines()
    i = 0
    title = ""
    subtitle = ""

    while i < len(lines):
        line = lines[i]
        if line.startswith("# ") and not title:
            title = strip_md(line[2:].strip())
            i += 1
            continue
        if not line.strip() and not title:
            i += 1
            continue
        if title and not subtitle and line.strip() and not line.startswith("#"):
            subtitle = strip_md(line.strip())
            i += 1
            continue
        if line.startswith("## "):
            blocks.append({"type": "section", "title": strip_md(line[3:].strip()), "items": []})
            i += 1
            continue
        if line.startswith("### "):
            if not blocks or blocks[-1]["type"] != "section":
                blocks.append({"type": "section", "title": strip_md(line[4:].strip()), "items": []})
            else:
                blocks[-1]["items"].append({"kind": "subheading", "text": strip_md(line[4:].strip())})
            i += 1
            continue
        if line.strip() == "---":
            i += 1
            continue
        img = re.match(r"^!\[[^\]]*\]\(([^)]+)\)\s*$", line.strip())
        if img:
            blocks.append({"type": "image", "path": img.group(1)})
            i += 1
            continue
        bullet = re.match(r"^[-*] (.+)$", line.strip())
        if bullet:
            if blocks and blocks[-1]["type"] == "section":
                blocks[-1]["items"].append({"kind": "bullet", "text": strip_md(bullet.group(1))})
            i += 1
            continue
        numbered = re.match(r"^(\d+)\. (.+)$", line.strip())
        if numbered:
            if blocks and blocks[-1]["type"] == "section":
                blocks[-1]["items"].append(
                    {
                        "kind": "numbered",
                        "num": numbered.group(1),
                        "text": strip_md(numbered.group(2)),
                    }
                )
            i += 1
            continue
        if line.strip():
            if blocks and blocks[-1]["type"] == "section":
                blocks[-1]["items"].append({"kind": "para", "text": strip_md(line.strip())})
        i += 1

    return [{"type": "title", "title": title, "subtitle": subtitle}, *blocks]


def resolve_image(path: str) -> Path | None:
    name = Path(path).name
    candidate = IMG_DIR / name
    return candidate if candidate.is_file() else None


def set_title_style(shape, size: int = 32, bold: bool = True) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.alignment = PP_ALIGN.LEFT


def add_text_slide(prs: Presentation, title: str, lines: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(MARGIN, Inches(0.35), CONTENT_W, Inches(0.9))
    title_box.text_frame.text = title
    set_title_style(title_box, size=28)

    body = slide.shapes.add_textbox(MARGIN, Inches(1.25), CONTENT_W, SLIDE_H - Inches(1.55))
    tf = body.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(16)
        p.space_after = Pt(6)
        p.level = 0


def add_image_slide(prs: Presentation, title: str, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(MARGIN, Inches(0.25), CONTENT_W, Inches(0.7))
    title_box.text_frame.text = title
    set_title_style(title_box, size=22)

    max_w = CONTENT_W
    max_h = SLIDE_H - Inches(1.2)
    slide.shapes.add_picture(str(image_path), MARGIN, Inches(0.95), width=max_w)
    pic = slide.shapes[-1]
    if pic.height > max_h:
        ratio = max_h / pic.height
        pic.height = int(max_h)
        pic.width = int(pic.width * ratio)
        pic.left = int((SLIDE_W - pic.width) / 2)


def chunk_lines(lines: list[str], max_lines: int = 10) -> list[list[str]]:
    if len(lines) <= max_lines:
        return [lines]
    chunks: list[list[str]] = []
    current: list[str] = []
    for line in lines:
        current.append(line)
        if len(current) >= max_lines:
            chunks.append(current)
            current = []
    if current:
        chunks.append(current)
    return chunks


def build_pptx(md_path: Path, out_path: Path) -> Path:
    blocks = parse_blocks(md_path.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_block = blocks[0]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(MARGIN, Inches(2.2), CONTENT_W, Inches(1.2))
    title_box.text_frame.text = title_block.get("title", "사용 가이드")
    set_title_style(title_box, size=40)
    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    sub = title_block.get("subtitle", "")
    if sub:
        sub_box = slide.shapes.add_textbox(MARGIN, Inches(3.5), CONTENT_W, Inches(1.0))
        sub_box.text_frame.text = sub
        set_title_style(sub_box, size=20, bold=False)
        sub_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    current_section = ""
    for block in blocks[1:]:
        if block["type"] == "section":
            current_section = block["title"]
            lines: list[str] = []
            for item in block["items"]:
                if item["kind"] == "subheading":
                    lines.append(item["text"])
                elif item["kind"] == "bullet":
                    lines.append(f"• {item['text']}")
                elif item["kind"] == "numbered":
                    lines.append(f"{item['num']}. {item['text']}")
                elif item["kind"] == "para":
                    lines.append(item["text"])
            for chunk in chunk_lines(lines):
                suffix = f" ({chunk[0][:12]}…)" if len(chunk_lines(lines)) > 1 else ""
                add_text_slide(prs, current_section + suffix, chunk)
        elif block["type"] == "image":
            img = resolve_image(block["path"])
            if img:
                add_image_slide(prs, current_section or "화면 안내", img)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def main() -> int:
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_OUT
    if not GUIDE_MD.is_file():
        print(f"가이드 파일 없음: {GUIDE_MD}", file=sys.stderr)
        return 1
    path = build_pptx(GUIDE_MD, out)
    print(path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
